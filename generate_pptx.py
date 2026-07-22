import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

prs = Presentation()
logo_path = r"C:\Users\ASUS\exam_java_vendredi\WhatsApp Image 2026-02-20 at 12.20.22.jpeg"

DARK_BLUE = RGBColor(10, 34, 64)
LIGHT_BLUE = RGBColor(0, 112, 192)
ORANGE = RGBColor(237, 125, 49)
WHITE = RGBColor(255, 255, 255)

def add_logo(slide, left=Inches(8.5), top=Inches(0.1), width=Inches(1.0)):
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, left, top, width=width)

def create_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Projet ERP LBP REFACTO"
    p.font.bold = True
    p.font.size = Pt(50)
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.add_paragraph()
    p2.text = "Gestion de projet avec le framework Scrum\nÉvaluation Continue - Exposé"
    p2.font.size = Pt(24)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    
    add_logo(slide, left=Inches(4.25), top=Inches(5.2), width=Inches(1.5))

def create_content_slide(title_text, content_pairs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    
    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    
    # Banner
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    banner.fill.solid()
    banner.fill.fore_color.rgb = DARK_BLUE
    banner.line.color.rgb = DARK_BLUE
    
    # Title Text
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(7.5), Inches(0.8))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    
    # Body
    txBox_body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf_body = txBox_body.text_frame
    tf_body.word_wrap = True
    
    for i, (heading, body) in enumerate(content_pairs):
        if heading:
            ph = tf_body.add_paragraph() if i > 0 else tf_body.paragraphs[0]
            ph.text = heading
            ph.font.bold = True
            ph.font.color.rgb = LIGHT_BLUE
            ph.font.size = Pt(20)
        
        if body:
            pb = tf_body.add_paragraph()
            pb.text = body
            pb.font.color.rgb = RGBColor(60, 60, 60)
            pb.font.size = Pt(16)
            pb.level = 1 if heading else 0
            
    add_logo(slide, left=Inches(8.8), top=Inches(0.1), width=Inches(1.0))
    return slide

create_title_slide()

create_content_slide("Sommaire de la présentation", [
    ("", "1. Présentation de l'Entreprise LBP"),
    ("", "2. Le Projet ERP LBP"),
    ("", "3. Le Framework Scrum : Définitions et Rôles"),
    ("", "4. Notre rythme de travail (Les Événements)"),
    ("", "5. Nos outils de suivi (Les Artefacts)"),
    ("", "6. Découpage fonctionnel (3 modules)"),
    ("", "7. Extrait de notre Product Backlog"),
    ("", "8. Conclusion et bilan")
])

create_content_slide("1. Présentation de l'entreprise LBP", [
    ("Qui est LBP ?", "LBP est une société spécialisée dans la logistique, l'export et l'import de fret."),
    ("Présence internationale :", "L'entreprise est implantée en France, au Sénégal et en Côte d'Ivoire, et compte un effectif de 43 personnes."),
    ("Problématique actuelle :", "Avec cette répartition géographique, la gestion RH et administrative (fichiers Excel, emails) devenait trop complexe et éparpillée."),
    ("L'objectif de notre projet :", "• Centraliser l'information pour gagner en efficacité.\n• Sécuriser l'accès aux données des collaborateurs (congés, rôles).\n• Se doter d'un outil sur mesure, évolutif et moderne.")
])

create_content_slide("2. Le Projet : ERP LBP REFACTO", [
    ("Notre solution pour LBP :", "Un ERP (Enterprise Resource Planning) modulaire."),
    ("La technique :", "• Architecture PHP MVC légère (Modèle-Vue-Contrôleur).\n• Socle technique : PHP 8, MySQL, Apache."),
    ("L'enjeu projet :", "• Construire ce système bloc par bloc en évitant de s'éparpiller.\n• Utiliser Scrum pour livrer de la valeur rapidement et s'adapter.")
])

create_content_slide("3. Scrum : Les Rôles et notre organisation", [
    ("L'Équipe Scrum :", "Définition : Pluridisciplinaire, auto-gérée et concentrée sur un Product Goal."),
    ("Product Owner (PO) :", "Théorie : Maximise la valeur du produit et gère le Backlog.\nChez nous : Il a joué le rôle de client pour LBP, priorisé le module RH et écrit nos User Stories."),
    ("Scrum Master (SM) :", "Théorie : Facilitateur et garant du cadre Scrum (Servant-Leader).\nChez nous : Il animait nos réunions, chronométrait, et aidait à résoudre les conflits Git."),
    ("Developers (Équipe) :", "Théorie : Les professionnels qui créent l'incrément à chaque Sprint.\nChez nous : L'équipe solidaire qui a codé l'architecture MVC et les pages PHP.")
])

create_content_slide("4. Scrum : Les Événements (Cérémonies)", [
    ("Le Sprint (1 mois max) :", "C'est le cœur de Scrum. Il contient 4 événements formels :"),
    ("1. Sprint Planning :", "Théorie : Définit le but (Sprint Goal) et comment le travail sera réalisé.\nChez nous : On sélectionnait ensemble les fonctionnalités à coder pour le Sprint."),
    ("2. Daily Scrum :", "Théorie : Point de 15 min de l'équipe pour se synchroniser.\nChez nous : Point rapide pour s'entraider sur le code (ex: 'Je bloque sur la session PHP')."),
    ("3. Sprint Review :", "Théorie : Présentation de l'Incrément pour inspection et adaptation.\nChez nous : On démontrait le module finalisé (comme la page de connexion)."),
    ("4. Sprint Retrospective :", "Théorie : Réflexion sur l'amélioration continue.\nChez nous : Moment d'honnêteté pour mieux s'organiser sur le Sprint suivant.")
])

create_content_slide("5. Scrum : Les Artefacts", [
    ("La Transparence :", "Les artefacts sont les outils qui apportent la transparence pour l'inspection."),
    ("Product Backlog :", "Théorie : Liste unique, ordonnée et émergente de tout ce qui est nécessaire.\nChez nous : Notre 'To-Do List' globale de l'ERP gérée exclusivement par le PO."),
    ("Sprint Backlog :", "Théorie : L'objectif du Sprint et les éléments du Backlog sélectionnés.\nChez nous : Notre plan d'action immédiat sur lequel on s'était engagé."),
    ("L'Incrément :", "Théorie : Résultat utilisable qui répond à la Definition of Done.\nChez nous : Une version de l'ERP LBP qui tourne localement, testée et sans erreurs critiques.")
])

create_content_slide("6. Découpage fonctionnel de l'ERP", [
    ("1. Le Portail et l'Authentification", "• Connexion sécurisée via Email/Mot de passe.\n• Gestion des sessions (expiration automatique).\n• Affichage dynamique du dashboard selon les droits."),
    ("2. Le Module Ressources Humaines (RH)", "• Création et édition d'une fiche employé dans la BDD.\n• Visualisation de l'annuaire du personnel.\n• Soumission et traitement des demandes de congés."),
    ("3. Le Module Administration", "• Attribution des rôles et des autorisations (Administrateur / Standard).\n• Activation du mode maintenance de l'ERP.\n• Consultation des logs d'erreurs système.")
])

# Slide Table
slide = prs.slides.add_slide(prs.slide_layouts[6])
banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
banner.fill.solid()
banner.fill.fore_color.rgb = DARK_BLUE
banner.line.color.rgb = DARK_BLUE
txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(7.5), Inches(0.8))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "7. Extrait de notre Product Backlog"
p.font.bold = True
p.font.size = Pt(32)
p.font.color.rgb = WHITE
add_logo(slide, left=Inches(8.8), top=Inches(0.1), width=Inches(1.0))

x, y, cx, cy = Inches(0.2), Inches(1.5), Inches(9.5), Inches(5)
shape = slide.shapes.add_table(10, 4, x, y, cx, cy)
table = shape.table

columns = ["Module", "User Story", "Priorité", "Taille"]
for i, col_name in enumerate(columns):
    cell = table.cell(0, i)
    cell.text = col_name
    cell.fill.solid()
    cell.fill.fore_color.rgb = LIGHT_BLUE
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = WHITE
            run.font.bold = True

data = [
    ("Portail", "En tant qu'utilisateur, je veux m'authentifier, afin d'accéder à l'ERP.", "1 - Critique", "5"),
    ("Portail", "En tant que système, je veux expirer la session, pour sécuriser l'accès.", "2 - Haute", "3"),
    ("Portail", "En tant qu'utilisateur, je veux voir le dashboard, pour choisir un module.", "1 - Critique", "5"),
    ("RH", "En tant que RH, je veux ajouter un employé, pour peupler la BDD.", "2 - Haute", "8"),
    ("RH", "En tant qu'employé, je veux voir l'annuaire, pour contacter un collègue.", "3 - Moyenne", "3"),
    ("RH", "En tant qu'employé, je veux demander un congé, pour mes vacances.", "3 - Moyenne", "8"),
    ("Admin", "En tant qu'admin, je veux attribuer des rôles, pour cloisonner les accès.", "2 - Haute", "5"),
    ("Admin", "En tant qu'admin, je veux activer la maintenance, pour faire des MAJ.", "4 - Basse", "3"),
    ("Admin", "En tant qu'admin, je veux visualiser les logs, pour tracer les bugs.", "4 - Basse", "2"),
]

for row_idx, row_data in enumerate(data):
    for col_idx, cell_data in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = cell_data
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(50, 50, 50)
                run.font.size = Pt(14)

table.columns[0].width = Inches(1.0)
table.columns[1].width = Inches(6.0)
table.columns[2].width = Inches(1.5)
table.columns[3].width = Inches(1.0)

# Slide Conclusion
create_content_slide("8. Conclusion et Bilan", [
    ("Ce que Scrum a apporté à notre projet ERP LBP :", ""),
    ("Une clarté théorique et pratique :", "Chaque cérémonie et chaque artefact a trouvé sa place dans notre travail."),
    ("Focus sur la valeur :", "En livrant des incréments par itération, nous avions très vite une application fonctionnelle au lieu d'un énorme code non testé à la fin."),
    ("Capacité d'adaptation :", "Si un module s'avérait trop complexe, on pouvait réestimer le Backlog avec le PO."),
    ("En résumé :", "Scrum a été un véritable pilier pour la réussite et l'organisation de ce projet d'étude.")
])

# Slide Fin
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = DARK_BLUE

txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "Merci de votre attention !"
p.font.bold = True
p.font.size = Pt(44)
p.font.color.rgb = ORANGE
p.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
tf2 = txBox2.text_frame
p2 = tf2.add_paragraph()
p2.text = "Avez-vous des questions sur notre organisation Scrum ou sur l'ERP LBP ?"
p2.font.size = Pt(24)
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.CENTER

add_logo(slide, left=Inches(4.25), top=Inches(5.0), width=Inches(1.5))


prs.save("Presentation_Scrum_ERP_LBP_V4_Design.pptx")
print("Presentation generated successfully!")
